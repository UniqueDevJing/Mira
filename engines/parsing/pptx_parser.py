"""PPT/PPTX 解析: python-pptx 抽取幻灯片标题/正文/表格 → UIRDocument。

依赖 python-pptx (lazy import: 未安装时仅 parse() 报错, 不影响模块导入与其他 parser)。
标题形状按占位符类型 (TITLE / CENTER_TITLE) 判定; 文本框正文作为段落; 表格结构化保留。
"""

import hashlib
import logging
from pathlib import Path

from engines.parsing.models import UIRDocument

logger = logging.getLogger(__name__)


class PptxParser:
    supported_types = (".pptx", ".ppt")

    def parse(self, file_path: str) -> UIRDocument:
        from pptx import Presentation

        prs = Presentation(file_path)
        blocks: list[dict] = []
        tables: list[dict] = []
        doc_id = hashlib.sha256(Path(file_path).read_bytes()).hexdigest()[:16]

        for si, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                    rows = [r for r in rows if any(r)]
                    if rows:
                        tables.append({"page_num": si, "bbox": [], "matrix": rows, "headers": rows[0]})
                        preview = " | ".join(rows[0]) + " ; " + " / ".join(
                            " | ".join(r) for r in rows[1 : min(len(rows), 9)]
                        )
                        blocks.append(
                            {"type": "paragraph", "bbox": [], "content": preview, "page_num": si, "metadata": {"is_table": True}}
                        )
                    continue
                txt = shape.text.strip() if shape.has_text_frame else ""
                if not txt:
                    continue
                is_title = False
                try:
                    pf = shape.placeholder_format
                    if pf is not None and pf.type is not None:
                        is_title = pf.type in (1, 13)  # TITLE / CENTER_TITLE
                except Exception:  # noqa: BLE001 — 占位符探测失败按非标题处理, 不阻断解析
                    is_title = False
                if is_title:
                    blocks.append(
                        {"type": "title", "bbox": [], "content": txt, "page_num": si, "metadata": {"heading_level": 1}}
                    )
                else:
                    blocks.append({"type": "paragraph", "bbox": [], "content": txt, "page_num": si, "metadata": {}})

        update_time = int(Path(file_path).stat().st_mtime)
        return UIRDocument(
            doc_id=doc_id,
            source={"type": "pptx", "path": file_path},
            pages=[{"page_num": 1, "blocks": blocks}],
            tables=tables,
            update_time=update_time,
        )
